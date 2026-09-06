#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
🎓 URY Engine — 멀티포맷 강의자료 파서 통합 모듈 v1.0 (doc_parser.py)
- 지원 포맷: .pdf, .pptx, .ppt, .hwpx, .hwp, .ipynb, .docx, .doc, .py, .sql, .txt, .md
- 주요 특징:
  1) 파워포인트(.pptx): 슬라이드 본문 + 교수님 발표자 노트(Notes) 자동 분리 추출
  2) 한글 문서(.hwpx, .hwp): 순수 파이썬 Zip/XML 및 바이너리 스트림 텍스트 파싱
  3) 주피터 노트북(.ipynb): 마크다운 설명과 코딩 셀 분리 파싱
  4) 소스코드(.py, .sql): 인코딩 자동 판별 및 주석 구조화 텍스트 반환
  5) 외부 툴/오피스 종속성 없는 Fallback 파싱 엔진 내장
"""

import os
import re
import json
import zipfile
import zlib
import unicodedata
import xml.etree.ElementTree as ET
from typing import Dict, Any, List


# --- 라이브러리 임포트 시도 (Fallback 보장) ---
try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False


def parse_document(file_path: str) -> Dict[str, Any]:
    """
    모든 형식의 강의자료 문서를 파싱하여 정제된 텍스트 및 구조 데이터를 반환하는 단일 인터페이스
    """
    if not os.path.exists(file_path):
        return {"status": "error", "error": f"파일이 존재하지 않습니다: {file_path}", "full_text": ""}

    file_name = os.path.basename(file_path)
    ext = os.path.splitext(file_name)[1].lower()

    result = {
        "status": "success",
        "file_path": file_path,
        "file_name": file_name,
        "ext": ext,
        "full_text": "",
        "sections": [],
        "notes_text": "",
        "metadata": {}
    }

    try:
        if ext == ".pdf":
            _parse_pdf(file_path, result)
        elif ext in [".pptx", ".ppt"]:
            _parse_pptx(file_path, result)
        elif ext == ".hwpx":
            _parse_hwpx(file_path, result)
        elif ext == ".hwp":
            _parse_hwp(file_path, result)
        elif ext == ".ipynb":
            _parse_ipynb(file_path, result)
        elif ext in [".docx", ".doc"]:
            _parse_docx(file_path, result)
        elif ext in [".html", ".htm"]:
            _parse_html(file_path, result)
        elif ext in [".py", ".sql", ".sh", ".c", ".cpp", ".java", ".js", ".css"]:
            _parse_code(file_path, result)
        elif ext in [".txt", ".md", ".csv", ".json"]:
            _parse_text(file_path, result)
        else:
            # 기본 텍스트 시도
            _parse_text(file_path, result)

    except Exception as e:
        result["status"] = "partial_error"
        result["error"] = str(e)
        if not result["full_text"]:
            result["full_text"] = f"[오류: {ext} 파일 파싱 중 예외 발생 - {str(e)}]"

    return result


def _parse_pdf(file_path: str, res: Dict[str, Any]):
    """PDF 파싱 (PyMuPDF -> pypdf fallback)"""
    full_text_list = []
    sections = []

    if HAS_FITZ:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                full_text_list.append(f"--- [Page {page_num}] ---\n{text}")
                sections.append({"page": page_num, "text": text})
        doc.close()
    elif HAS_PYPDF:
        reader = pypdf.PdfReader(file_path)
        for page_num, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                full_text_list.append(f"--- [Page {page_num}] ---\n{text}")
                sections.append({"page": page_num, "text": text})
    else:
        # 텍스트로 무작정 시도
        _parse_text(file_path, res)
        return

    res["full_text"] = "\n\n".join(full_text_list)
    res["sections"] = sections


def _parse_pptx(file_path: str, res: Dict[str, Any]):
    """PPTX 파싱: 슬라이드 텍스트 + 교수님 발표자 노트(Notes) 분리 추출"""
    full_text_list = []
    sections = []
    all_notes = []

    if HAS_PPTX and file_path.endswith(".pptx"):
        prs = pptx.Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            notes_text = ""

            # 본문 도형 텍스트 추출
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)

            # 발표자 노트 추출
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            body_text = "\n".join(slide_text_parts)
            section_entry = {
                "slide_num": i,
                "title": f"Slide {i}",
                "content": body_text,
                "notes": notes_text
            }
            sections.append(section_entry)

            combined_str = f"[Slide {i}]\n{body_text}"
            if notes_text:
                combined_str += f"\n🎙️ [교수님 발표자 노트/Notes]: {notes_text}"
                all_notes.append(f"Slide {i}: {notes_text}")

            full_text_list.append(combined_str)

    else:
        # Zip XML Fallback 파싱
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, 'r') as z:
                slide_files = sorted([f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])
                for i, s_file in enumerate(slide_files, 1):
                    xml_content = z.read(s_file).decode('utf-8', errors='ignore')
                    try:
                        root = ET.fromstring(xml_content)
                        texts = [elem.text for elem in root.iter() if elem.text and elem.text.strip()]
                    except Exception:
                        texts = re.findall(r'<a:t[^>]*>(.*?)</a:t>', xml_content)
                    body_text = " ".join(texts)
                    sections.append({"slide_num": i, "content": body_text, "notes": ""})
                    full_text_list.append(f"[Slide {i}]\n{body_text}")


    res["full_text"] = "\n\n".join(full_text_list)
    res["sections"] = sections
    res["notes_text"] = "\n".join(all_notes)


def _parse_hwpx(file_path: str, res: Dict[str, Any]):
    """.hwpx (한글 XML 표준) Zip 파싱"""
    if not zipfile.is_zipfile(file_path):
        _parse_text(file_path, res)
        return

    full_text_parts = []
    with zipfile.ZipFile(file_path, 'r') as z:
        section_files = sorted([f for f in z.namelist() if 'Contents/section' in f and f.endswith('.xml')])
        for s_file in section_files:
            xml_data = z.read(s_file).decode('utf-8', errors='ignore')
            try:
                root = ET.fromstring(xml_data)
                for elem in root.iter():
                    if elem.tag.endswith('}t') or elem.tag.endswith(':t') or elem.tag == 't':
                        if elem.text:
                            full_text_parts.append(elem.text.strip())
            except Exception:
                # XML 네임스페이스 오류 시 정규식 fallback
                text_matches = re.findall(r'<(?:hp:)?t[^>]*>(.*?)</(?:hp:)?t>', xml_data, re.DOTALL)
                for tm in text_matches:
                    clean = re.sub(r'<[^>]+>', '', tm).strip()
                    if clean:
                        full_text_parts.append(clean)

    text_content = "\n".join([t for t in full_text_parts if t])
    res["full_text"] = text_content



def _parse_hwp(file_path: str, res: Dict[str, Any]):
    """.hwp 5.0 바이너리 스트림 파싱 / 텍스트 추출 fallback"""
    full_text_parts = []
    
    try:
        with open(file_path, 'rb') as f:
            data = f.read()

        raw_text = data.decode('utf-16le', errors='ignore')
        clean_text = re.sub(r'[^\w\s\n\.\,\:\;\(\)\[\]\-\+\=\/\%\@\!\?가-힣A-Za-z0-9]', ' ', raw_text)
        clean_lines = [line.strip() for line in clean_text.splitlines() if len(line.strip()) > 3]
        
        if len(clean_lines) > 5:
            full_text_parts = clean_lines
        else:
            raw_cp = data.decode('cp949', errors='ignore')
            clean_cp = re.sub(r'[^\w\s\n\.\,\:\;\(\)\[\]\-\+\=\/\%\@\!\?가-힣A-Za-z0-9]', ' ', raw_cp)
            full_text_parts = [line.strip() for line in clean_cp.splitlines() if len(line.strip()) > 3]

    except Exception:
        pass

    res["full_text"] = "\n".join(full_text_parts[:500])


def _parse_ipynb(file_path: str, res: Dict[str, Any]):
    """주피터 노트북 파싱 (.ipynb)"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        nb_data = json.load(f)

    full_text_list = []
    sections = []

    cells = nb_data.get("cells", [])
    for idx, cell in enumerate(cells, 1):
        cell_type = cell.get("cell_type", "unknown")
        source = cell.get("source", [])
        if isinstance(source, list):
            cell_text = "".join(source).strip()
        else:
            cell_text = str(source).strip()

        if not cell_text:
            continue

        if cell_type == "markdown":
            formatted = f"### 📌 [Markdown Cell {idx}]\n{cell_text}"
        elif cell_type == "code":
            formatted = f"```python\n# [Code Cell {idx}]\n{cell_text}\n```"
        else:
            formatted = f"[{cell_type.upper()} Cell {idx}]\n{cell_text}"

        full_text_list.append(formatted)
        sections.append({"cell_num": idx, "cell_type": cell_type, "content": cell_text})

    res["full_text"] = "\n\n".join(full_text_list)
    res["sections"] = sections


def _parse_docx(file_path: str, res: Dict[str, Any]):
    """Word 문서 파싱 (.docx)"""
    if HAS_DOCX and file_path.endswith(".docx"):
        doc = docx.Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        res["full_text"] = "\n\n".join(paragraphs)
    else:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, 'r') as z:
                xml_content = z.read('word/document.xml')
                root = ET.fromstring(xml_content)
                texts = [elem.text for elem in root.iter() if elem.text and elem.text.strip()]
                res["full_text"] = "\n".join(texts)
        else:
            _parse_text(file_path, res)


def _parse_html(file_path: str, res: Dict[str, Any]):
    """HTML 웹/문서 파일 (.html, .htm) 파싱 (태그 제거 및 구조적 텍스트 추출)"""
    import html.parser
    raw = _read_file_text_with_encodings(file_path)

    class _HTMLTextExtractor(html.parser.HTMLParser):
        def __init__(self):
            super().__init__()
            self.pieces = []
            self.skip = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style", "head", "meta", "noscript"):
                self.skip = True
            elif tag in ("p", "div", "h1", "h2", "h3", "h4", "h5", "h6", "tr", "li", "br"):
                self.pieces.append("\n")

        def handle_endtag(self, tag):
            if tag in ("script", "style", "head", "meta", "noscript"):
                self.skip = False
            elif tag in ("p", "div", "h1", "h2", "h3", "h4", "h5", "h6", "tr", "table"):
                self.pieces.append("\n")

        def handle_data(self, data):
            if not self.skip:
                clean = data.strip()
                if clean:
                    self.pieces.append(clean + " ")

    extractor = _HTMLTextExtractor()
    try:
        extractor.feed(raw)
        lines = [line.strip() for line in "".join(extractor.pieces).splitlines() if line.strip()]
        res["full_text"] = "\n".join(lines)
    except Exception:
        import re
        clean = re.sub(r"<(script|style).*?>.*?</\1>", "", raw, flags=re.DOTALL | re.IGNORECASE)
        clean = re.sub(r"<[^>]+>", " ", clean)
        lines = [line.strip() for line in clean.splitlines() if line.strip()]
        res["full_text"] = "\n".join(lines)


def _parse_code(file_path: str, res: Dict[str, Any]):
    """소스 코드 파일 (.py, .sql 등) 파싱"""
    content = _read_file_text_with_encodings(file_path)
    file_name = os.path.basename(file_path)
    ext = os.path.splitext(file_name)[1].lstrip('.')
    
    res["full_text"] = f"```{ext}\n# [소스코드 파일: {file_name}]\n{content}\n```"


def _parse_text(file_path: str, res: Dict[str, Any]):
    """일반 텍스트 파일 파싱"""
    content = _read_file_text_with_encodings(file_path)
    res["full_text"] = content


def _read_file_text_with_encodings(file_path: str) -> str:
    """다양한 텍스트 인코딩 순차 시도 (utf-8, cp949, euc-kr, latin-1)"""
    encodings = ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            return f"[텍스트 읽기 오류: {str(e)}]"
    
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        print(f"🔍 파싱 테스트: {test_file}")
        parsed = parse_document(test_file)
        print(f"상태: {parsed['status']}")
        print(f"확장자: {parsed['ext']}")
        print(f"추출 텍스트 길이: {len(parsed['full_text'])} 자")
        print("--- 상위 300자 출력 ---")
        print(parsed['full_text'][:300])
